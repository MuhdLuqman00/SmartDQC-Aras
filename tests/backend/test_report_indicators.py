"""Tests for Feature #15 — KKM-formatted report generation (PDF + PPTX)."""
from io import BytesIO

import pdfplumber
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

from backend.export.report import build_pptx_bytes, build_pdf_bytes
from backend.export.report_template_spec import (
    KKM_TEAL, SECTION_LABELS,
)
from backend.export.charts import (
    chart_quality_bar, chart_nutritional_rates, chart_kpi_vs_target,
)
import pandas as pd
from backend.eda.kpi import compute_kpi_dashboard


def _pdf_text(data: bytes) -> str:
    """Extract all text from a PDF using pdfplumber."""
    with pdfplumber.open(BytesIO(data)) as pdf:
        return "\n".join(
            page.extract_text() or "" for page in pdf.pages
        )


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

def _eda():
    """Mirrors run_eda()'s ACTUAL output schema.

    The old fixture used legacy keys (summary / quality.overall_score /
    flat indicators / outliers.total_flagged) that run_eda never emits —
    the same fiction that made the report's quality overview render blank
    in production. This pins the real producer -> consumer contract.
    """
    return {
        "source_type":   "myvass",
        "total_rows":    100,
        "total_columns": 12,
        "data_quality_score": {
            "score": 82.0, "grade": "B", "label": "Baik",
            "breakdown": {
                "field_coverage":   {"score": 16.5, "max": 20},
                "ic_validity":      {"score": 15.0, "max": 15},
                "missing_critical": {"score": 20.0, "max": 20},
                "duplicates":       {"score": 12.0, "max": 15},
                "bmi_consistency":  {"score": 10.0, "max": 10},
                "spelling":         {"score":  4.5, "max":  5},
                "zscore_coverage":  {"score":  4.0, "max": 15},
            },
        },
        "data_completeness": {"pct_complete": 91.5,
                              "pct_missing_critical": 0.08},
        "outliers": {
            "berat_kg": {"column": "berat_kg", "combined_outliers": 7},
        },
        "indicators": {
            "bawah_5_tahun": {
                "stunting": {
                    "label": "Bantut (HAZ < -2)",
                    "overall": {"pct": 18.5, "n_affected": 18,
                                "n_total": 100},
                },
                "wasting": {
                    "label": "Susut (BAZ < -2)",
                    "overall": {"pct": 3.2, "n_affected": 3,
                                "n_total": 100},
                },
            },
        },
    }


def _narrative():
    return {
        "executive_summary": {
            "bm": "Ringkasan ujian dalam Bahasa Malaysia.",
            "en": "Test executive summary in English.",
        },
        "recommendations": [
            {"action": "Increase monitoring", "priority": "HIGH",
             "bm": "Tingkatkan pemantauan.", "en": "Increase monitoring."},
        ],
    }


def _kpi():
    # Build from the REAL producer so the report consumers are pinned to the
    # actual compute_kpi_dashboard contract (indicators + by_state/by_daerah
    # with nested `rates`), not a hand-written fiction that drifts from it.
    df = pd.DataFrame({
        "stunting":    [1] * 18 + [0] * 82,
        "wasting":     [1] * 3  + [0] * 97,
        "underweight": [1] * 9  + [0] * 91,
        "overweight":  [1] * 8  + [0] * 92,
        "NEGERI":      ["Selangor"] * 50 + ["Johor"] * 50,
    })
    return compute_kpi_dashboard(df)


# ---------------------------------------------------------------------------
# PPTX tests
# ---------------------------------------------------------------------------

def test_pptx_without_kpi_has_five_slides():
    """Cover + Exec Summary + Quality + Recommendations + Methodology = 5."""
    data = build_pptx_bytes(_eda(), _narrative())
    prs  = Presentation(BytesIO(data))
    assert len(prs.slides) == 5


def test_pptx_with_kpi_has_seven_slides():
    """Adds Indicator Table + Indicator Charts slides when kpi_result is provided
    (5 base slides + 2 KPI slides = 7)."""
    data = build_pptx_bytes(_eda(), _narrative(), kpi_result=_kpi())
    prs  = Presentation(BytesIO(data))
    assert len(prs.slides) == 7


def test_pptx_cover_slide_has_text():
    data = build_pptx_bytes(_eda(), _narrative(), district="Petaling")
    prs  = Presentation(BytesIO(data))
    cover_text = " ".join(
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
    )
    assert "SmartDQC" in cover_text
    assert "Petaling" in cover_text


def test_pptx_section_bars_use_brand_color():
    """Every non-cover slide should carry at least one section bar filled with
    the KKM brand color (navy after the v3 reskin; tracked via KKM_TEAL)."""
    data = build_pptx_bytes(_eda(), _narrative(), kpi_result=_kpi())
    prs  = Presentation(BytesIO(data))
    brand = RGBColor.from_string(KKM_TEAL.lstrip("#"))

    def _solid_rgb(sh):
        """fore_color.rgb for solid-filled shapes only; None otherwise."""
        try:
            if sh.fill.type == MSO_FILL.SOLID:
                return sh.fill.fore_color.rgb
        except (TypeError, AttributeError):
            pass
        return None

    for idx, slide in enumerate(prs.slides):
        if idx == 0:          # cover slide has dark bg, not a section bar
            continue
        brand_found = any(
            _solid_rgb(sh) == brand
            for sh in slide.shapes
            if sh.has_text_frame
        )
        assert brand_found, f"No brand-color section bar found on slide index {idx}"


def test_pptx_exec_summary_bilingual():
    data = build_pptx_bytes(_eda(), _narrative())
    prs  = Presentation(BytesIO(data))
    exec_slide_text = " ".join(
        sh.text_frame.text
        for sh in prs.slides[1].shapes
        if sh.has_text_frame
    )
    assert "Bahasa Malaysia" in exec_slide_text
    assert "English"         in exec_slide_text


def test_pptx_footer_contains_kkm():
    data = build_pptx_bytes(_eda(), _narrative())
    prs  = Presentation(BytesIO(data))
    all_text = " ".join(
        sh.text_frame.text
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
    )
    assert "Kementerian Kesihatan Malaysia" in all_text


def test_pptx_district_param_appears_in_output():
    data = build_pptx_bytes(_eda(), _narrative(), district="Klang Utara")
    prs  = Presentation(BytesIO(data))
    all_text = " ".join(
        sh.text_frame.text
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
    )
    assert "Klang Utara" in all_text


# ---------------------------------------------------------------------------
# PDF tests
# ---------------------------------------------------------------------------

def test_pdf_is_valid():
    data = build_pdf_bytes(_eda(), _narrative())
    assert data[:4] == b"%PDF"


def test_pdf_with_kpi_is_valid():
    data = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi())
    assert data[:4] == b"%PDF"


def test_pdf_contains_bilingual_section_labels():
    """BM and EN section labels must appear in extracted PDF text."""
    data = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi())
    text = _pdf_text(data)
    assert "RINGKASAN EKSEKUTIF"              in text
    assert "EXECUTIVE SUMMARY"               in text
    assert "INDICATOR TABLE BY DISTRICT"     in text
    assert "JADUAL PETUNJUK MENGIKUT DAERAH" in text
    assert "METHODOLOGY APPENDIX"            in text
    assert "LAMPIRAN METODOLOGI"             in text


def test_pdf_footer_contains_kkm():
    data = build_pdf_bytes(_eda(), _narrative(), district="Selangor")
    text = _pdf_text(data)
    assert "Kementerian Kesihatan Malaysia" in text
    assert "Selangor"                       in text


def test_pdf_quality_section_present():
    data = build_pdf_bytes(_eda(), _narrative())
    text = _pdf_text(data)
    assert "82"           in text   # overall_score from fixture
    assert "Completeness" in text


def test_pdf_recommendations_present():
    data = build_pdf_bytes(_eda(), _narrative())
    text = _pdf_text(data)
    assert "Increase monitoring" in text


def test_pdf_lang_bm_uses_bm_headers():
    """lang='bm' should use Malay column headers in indicator table."""
    data = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi(), lang="bm")
    text = _pdf_text(data)
    assert "Daerah" in text


# ---------------------------------------------------------------------------
# Chart unit tests
# ---------------------------------------------------------------------------

def test_chart_quality_bar_returns_png():
    png = chart_quality_bar(_eda())
    assert png is not None
    assert png[:8] == b"\x89PNG\r\n\x1a\n"


def test_chart_quality_bar_returns_none_without_quality_data():
    assert chart_quality_bar({}) is None


def test_chart_nutritional_rates_returns_png():
    png = chart_nutritional_rates(_kpi())
    assert png is not None
    assert png[:8] == b"\x89PNG\r\n\x1a\n"


def test_chart_nutritional_rates_returns_none_without_breakdown():
    assert chart_nutritional_rates({"by_state": [], "by_daerah": []}) is None
    assert chart_nutritional_rates({}) is None
    assert chart_nutritional_rates(None) is None


def test_chart_kpi_vs_target_returns_png():
    png = chart_kpi_vs_target(_kpi())
    assert png is not None
    assert png[:8] == b"\x89PNG\r\n\x1a\n"


def test_chart_kpi_vs_target_returns_none_without_kpis():
    assert chart_kpi_vs_target({}) is None
    assert chart_kpi_vs_target(None) is None


def test_pdf_with_kpi_embeds_charts():
    """Charts produce a larger PDF than the table-only version."""
    data_no_kpi  = build_pdf_bytes(_eda(), _narrative())
    data_with_kpi = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi())
    assert len(data_with_kpi) > len(data_no_kpi)


def test_pdf_explicit_no_charts_is_smaller_than_default():
    """An explicit empty chart set embeds no charts; default (None) embeds all."""
    full = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi(), charts=None)
    none = build_pdf_bytes(_eda(), _narrative(), kpi_result=_kpi(), charts=set())
    assert len(none) < len(full)


def test_pptx_quality_slide_with_chart_is_larger():
    """Quality chart image adds bytes to the PPTX."""
    data = build_pptx_bytes(_eda(), _narrative())
    assert len(data) > 10_000
