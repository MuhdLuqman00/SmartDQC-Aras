"""SmartDQC report builders — KKM-formatted PDF and PPTX output.

Format reference: MOH Malaysia Annual Report 2024 (teal section bars,
bilingual headers, alternating-row tables, footer stamp).

Content (Feature #15):
  1. Cover
  2. Executive Summary  (BM + EN — reused from Feature #9 narrative)
  3. Data Quality Overview
  4. Recommendations    (reused from Feature #9 narrative)
  5. Indicator Tables by District  (requires kpi_result)
  6. Methodology Appendix
"""
from __future__ import annotations

from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import A4
from reportlab.lib import colors as rl_colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, Image,
)

from backend.export.charts import chart_quality_bar, chart_nutritional_rates, chart_kpi_vs_target

from backend.export.report_template_spec import (
    KKM_TEAL, KKM_TEAL_DARK, KKM_TEAL_LIGHT,
    KKM_NAVY, KKM_MID_GRAY, KKM_RULE_LINE,
    SECTION_LABELS, FOOTER_TEMPLATE, METHODOLOGY_LINES,
    NUTRITIONAL_TABLE_HEADERS, KPI_TABLE_HEADERS,
)

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _sec(key: str, lang: str = "en") -> str:
    return SECTION_LABELS.get(key, {}).get(lang, key.upper())


def _hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _fmt(val) -> str:
    if val is None:
        return "-"
    try:
        return f"{float(val):.1f}"
    except (TypeError, ValueError):
        return str(val)


# ===========================================================================
# PPTX
# ===========================================================================

def _rgb(h: str) -> RGBColor:
    r, g, b = _hex_to_rgb(h)
    return RGBColor(r, g, b)


def _bg(slide, h: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(h)


def _box(slide, text, l, t, w, h, size=11, bold=False,
         color="#FFFFFF", align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _section_bar_pptx(slide, key: str):
    """Teal band across the top with 'EN TITLE  /  BM TITLE'."""
    label = f"{_sec(key, 'en')}  /  {_sec(key, 'bm')}"
    tb = slide.shapes.add_textbox(
        Inches(0), Inches(0), Inches(13.33), Inches(0.6),
    )
    tf = tb.text_frame
    tb.fill.solid()
    tb.fill.fore_color.rgb = _rgb(KKM_TEAL)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {label}"
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = _rgb("#FFFFFF")


def _footer_bar_pptx(slide, district: str):
    footer = FOOTER_TEMPLATE.format(district=district, year=date.today().year)
    tb = slide.shapes.add_textbox(
        Inches(0), Inches(7.2), Inches(13.33), Inches(0.28),
    )
    tf = tb.text_frame
    tb.fill.solid()
    tb.fill.fore_color.rgb = _rgb(KKM_TEAL_DARK)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = footer
    run.font.size = Pt(7.5)
    run.font.color.rgb = _rgb(KKM_TEAL_LIGHT)


def _pptx_table(slide, rows: list[list[str]], l, t, w, h,
                status_cols: list[int] | None = None):
    """Add a styled KKM table. status_cols = col indices to colour-code."""
    n_r, n_c = len(rows), len(rows[0])
    tbl = slide.shapes.add_table(
        n_r, n_c, Inches(l), Inches(t), Inches(w), Inches(h),
    ).table

    status_map = {"Green": KKM_TEAL, "Amber": "#E8A020", "Red": "#C0392B"}

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(10 if ri > 0 else 11)
            run.font.bold = ri == 0

            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(KKM_TEAL)
                run.font.color.rgb = _rgb("#FFFFFF")
            elif status_cols and ci in status_cols and val in status_map:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(status_map[val])
                run.font.color.rgb = _rgb("#FFFFFF")
                run.font.bold = True
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(KKM_TEAL_LIGHT)
                run.font.color.rgb = _rgb(KKM_NAVY)
            else:
                run.font.color.rgb = _rgb(KKM_NAVY)


# --- individual slides -------------------------------------------------------

def _slide_cover(prs, layout, eda: dict, district: str, date_range: str):
    s = prs.slides.add_slide(layout)
    _bg(s, KKM_TEAL_DARK)
    today  = date.today().strftime("%d %B %Y")
    # run_eda emits source_type/total_rows at the top level; the legacy
    # nested "summary" shape is kept only as a fallback.
    source = (eda.get("source_type")
              or eda.get("summary", {}).get("source_type", "N/A")).upper()
    rows   = eda.get("total_rows",
                     eda.get("summary", {}).get("total_rows", "N/A"))

    _box(s, "Kementerian Kesihatan Malaysia", 0.6, 0.5, 12, 0.45,
         size=11, color=KKM_TEAL_LIGHT)
    _box(s, "LAPORAN PEMAKANAN SmartDQC", 0.6, 1.1, 12, 1.1,
         size=34, bold=True)
    _box(s, "SMARTDQC NUTRITION REPORT", 0.6, 2.1, 12, 0.55,
         size=15, color=KKM_TEAL_LIGHT)

    meta = f"Daerah / District: {district}  |  Sumber / Source: {source}  |  Rekod / Records: {rows}"
    if date_range:
        meta += f"  |  Tempoh / Period: {date_range}"
    _box(s, meta, 0.6, 3.0, 12.0, 0.45, size=10, color=KKM_TEAL_LIGHT)
    _box(s, today, 0.6, 3.45, 12.0, 0.4, size=10, color="#FFFFFF")
    _footer_bar_pptx(s, district)


def _slide_exec_summary(prs, layout, narrative: dict, district: str):
    s = prs.slides.add_slide(layout)
    _bg(s, "#FFFFFF")
    _section_bar_pptx(s, "executive_summary")
    exec_sum = narrative.get("executive_summary", {})
    _box(s, "Bahasa Malaysia", 0.5, 0.75, 6.0, 0.38,
         size=10, bold=True, color=KKM_TEAL)
    _box(s, exec_sum.get("bm", "-"), 0.5, 1.15, 5.9, 5.7,
         size=10, color=KKM_NAVY)
    _box(s, "English", 7.0, 0.75, 6.0, 0.38,
         size=10, bold=True, color=KKM_TEAL)
    _box(s, exec_sum.get("en", "-"), 7.0, 1.15, 5.9, 5.7,
         size=10, color=KKM_NAVY)
    _footer_bar_pptx(s, district)


def _quality_overview_rows(eda: dict) -> list[list[str]]:
    """Quality-overview metric rows from run_eda()'s ACTUAL schema.

    The report historically read eda["quality"]["overall_score"] /
    "overall_completeness" / "missing_rate" and outliers["total_flagged"]
    — keys run_eda never emits. run_eda produces data_quality_score
    (score/grade/label), data_completeness (pct_complete /
    pct_missing_critical) and a per-column outliers dict. Reading the
    legacy keys made every value render as "-".
    """
    dq   = eda.get("data_quality_score") or {}
    dc   = eda.get("data_completeness") or {}
    outl = eda.get("outliers") or {}

    score = dq.get("score")
    grade = dq.get("grade")
    label = dq.get("label")
    if score is not None and grade:
        score_txt = f"{score} (Grade {grade}" + (f" - {label})" if label else ")")
    elif score is not None:
        score_txt = str(score)
    else:
        score_txt = "-"

    total_outliers = sum(
        int(v.get("combined_outliers", 0) or 0)
        for v in outl.values()
        if isinstance(v, dict)
    ) if isinstance(outl, dict) else 0

    pct_complete = dc.get("pct_complete")
    pct_missing  = dc.get("pct_missing_critical")
    return [
        ["Overall Quality Score / Skor Kualiti", score_txt],
        ["Completeness / Kelengkapan",
         f"{pct_complete}%" if pct_complete is not None else "-"],
        ["Missing Data Rate / Kadar Data Hilang",
         f"{pct_missing}%" if pct_missing is not None else "-"],
        ["Outliers Flagged / Pencilan Ditanda", str(total_outliers)],
    ]


def _slide_quality(prs, layout, eda: dict, district: str, charts: set[str] | None = None):
    s = prs.slides.add_slide(layout)
    _bg(s, KKM_TEAL_LIGHT)
    _section_bar_pptx(s, "quality_overview")

    rows = [["Metric / Metrik", "Value / Nilai"]]
    rows.extend(_quality_overview_rows(eda))

    _pptx_table(s, rows, l=0.4, t=0.75, w=6.0, h=min(0.45 * len(rows), 6.3))

    if charts is None or "quality_bar" in charts:
        chart_png = chart_quality_bar(eda)
        if chart_png:
            s.shapes.add_picture(BytesIO(chart_png), Inches(6.8), Inches(0.75), Inches(6.1), Inches(5.8))

    _footer_bar_pptx(s, district)


def _action_for(rec: dict, lang: str) -> str:
    """Pick the recommendation action title in the requested language;
    fall back to the other language, then to the legacy single-string
    `action` so cached narratives keep rendering."""
    if lang == "bm":
        return rec.get("action_bm") or rec.get("action_en") or rec.get("action") or ""
    return rec.get("action_en") or rec.get("action_bm") or rec.get("action") or ""


def _slide_recommendations(prs, layout, narrative: dict, district: str):
    s = prs.slides.add_slide(layout)
    _bg(s, "#FFFFFF")
    _section_bar_pptx(s, "recommendations")
    for i, rec in enumerate(narrative.get("recommendations", [])[:3]):
        y = 0.75 + i * 2.0
        priority = rec.get("priority", "").upper()
        action_bm = _action_for(rec, "bm")
        action_en = _action_for(rec, "en")
        # Priority pill spans both columns; per-language action sits with its body.
        _box(s, f"[{priority}]", 0.5, y, 1.2, 0.45,
             size=12, bold=True, color=KKM_NAVY)
        _box(s, action_bm, 1.8, y, 4.7, 0.45,
             size=12, bold=True, color=KKM_NAVY)
        _box(s, action_en, 6.7, y, 6.0, 0.45,
             size=12, bold=True, color=KKM_NAVY)
        _box(s, rec.get("bm", ""), 0.5, y + 0.48, 6.0, 1.3, size=10, color=KKM_MID_GRAY)
        _box(s, rec.get("en",  ""), 6.7, y + 0.48, 6.0, 1.3, size=10, color=KKM_MID_GRAY)
    _footer_bar_pptx(s, district)


def _slide_indicator_table(prs, layout, kpi_result: dict, district: str, lang: str = "en",
                           charts: set[str] | None = None):
    breakdown = kpi_result.get("by_daerah") or kpi_result.get("by_state") or []
    s = prs.slides.add_slide(layout)
    _bg(s, KKM_TEAL_LIGHT)
    _section_bar_pptx(s, "indicator_table")

    headers = NUTRITIONAL_TABLE_HEADERS[lang]
    rows = [headers]
    for row in breakdown[:8]:
        rates = row.get("rates") or {}
        rows.append([
            str(row.get("district") or row.get("state") or ""),
            str(row.get("n", "")),
            _fmt(rates.get("stunting")),
            _fmt(rates.get("wasting")),
            _fmt(rates.get("underweight")),
            _fmt(rates.get("overweight")),
        ])

    table_h = min(0.42 * len(rows), 3.5)
    if len(rows) > 1:
        _pptx_table(s, rows, l=0.4, t=0.75, w=12.5, h=table_h)
    else:
        _box(s, "No district breakdown available.", 0.5, 1.5, 12, 1.0,
             size=11, color=KKM_NAVY)

    # Honour the per-chart filter; default (charts is None) embeds both.
    want_rates  = charts is None or "nutritional_rates" in charts
    want_target = charts is None or "kpi_vs_target" in charts

    chart_y = 0.75 + table_h + 0.15
    chart_h = min(6.9 - chart_y, 3.2)
    if chart_h > 1.0:
        chart2 = chart_nutritional_rates(kpi_result) if want_rates else None
        chart3 = chart_kpi_vs_target(kpi_result) if want_target else None
        if chart2:
            s.shapes.add_picture(BytesIO(chart2), Inches(0.4), Inches(chart_y),
                                 Inches(6.4), Inches(chart_h))
        if chart3:
            s.shapes.add_picture(BytesIO(chart3), Inches(7.0), Inches(chart_y),
                                 Inches(6.0), Inches(chart_h))

    _footer_bar_pptx(s, district)


def _slide_methodology(prs, layout, district: str):
    s = prs.slides.add_slide(layout)
    _bg(s, KKM_TEAL_DARK)
    _section_bar_pptx(s, "methodology")
    body = "\n".join(f"  - {line}" for line in METHODOLOGY_LINES)
    _box(s, body, 0.5, 0.75, 12.3, 6.3, size=9.5)
    _footer_bar_pptx(s, district)


# --- public entry point ------------------------------------------------------

def build_pptx_bytes(
    eda_result: dict,
    narrative: dict,
    kpi_result: dict | None = None,
    district: str = "Malaysia",
    date_range: str = "",
    lang: str = "en",
    charts: set[str] | None = None,
) -> bytes:
    """Render the PPTX report.

    `charts` controls which chart PNGs are embedded:
      None         → default (every recommended chart)
      set of keys  → only those keys; known: quality_bar, nutritional_rates,
                     kpi_vs_target.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _slide_cover(prs, blank, eda_result, district, date_range)
    _slide_exec_summary(prs, blank, narrative, district)
    _slide_quality(prs, blank, eda_result, district, charts=charts)
    _slide_recommendations(prs, blank, narrative, district)

    if kpi_result:
        _slide_indicator_table(prs, blank, kpi_result, district, lang, charts=charts)

    _slide_methodology(prs, blank, district)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===========================================================================
# PDF
# ===========================================================================

def _pdf_styles():
    base = getSampleStyleSheet()
    cover_h = ParagraphStyle(
        "CoverH", parent=base["Normal"],
        fontSize=28, leading=34, fontName="Helvetica-Bold",
        textColor=rl_colors.white,
    )
    sec_label = ParagraphStyle(
        "SecLabel", parent=base["Normal"],
        fontSize=11, leading=15, fontName="Helvetica-Bold",
        textColor=rl_colors.white,
    )
    h2 = ParagraphStyle(
        "H2", parent=base["Normal"],
        fontSize=12, leading=16, fontName="Helvetica-Bold",
        textColor=rl_colors.HexColor(KKM_TEAL), spaceBefore=8,
    )
    body = ParagraphStyle(
        "Body", parent=base["Normal"],
        fontSize=10, leading=15,
        textColor=rl_colors.HexColor(KKM_NAVY),
    )
    small = ParagraphStyle(
        "Small", parent=base["Normal"],
        fontSize=8.5, leading=13,
        textColor=rl_colors.HexColor(KKM_MID_GRAY),
    )
    return cover_h, sec_label, h2, body, small


def _section_bar_pdf(label_en: str, label_bm: str, sec_label_style) -> Table:
    """Teal header bar matching KKM section style."""
    tbl = Table(
        [[Paragraph(f"{label_en}  /  {label_bm}", sec_label_style)]],
        colWidths=[17 * cm],
    )
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), rl_colors.HexColor(KKM_TEAL)),
        ("TOPPADDING",    (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("LEFTPADDING",   (0, 0), (-1, -1), 10),
    ]))
    return tbl


def _base_table_style() -> list:
    return [
        ("BACKGROUND",    (0, 0), (-1, 0),  rl_colors.HexColor(KKM_TEAL)),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  rl_colors.white),
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTSIZE",      (0, 0), (-1, -1), 9),
        ("GRID",          (0, 0), (-1, -1), 0.4, rl_colors.HexColor(KKM_RULE_LINE)),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1),
         [rl_colors.white, rl_colors.HexColor(KKM_TEAL_LIGHT)]),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING",    (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING",   (0, 0), (-1, -1), 5),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 5),
    ]


def _make_footer_canvas(district: str, year: int):
    from reportlab.pdfgen.canvas import Canvas

    class KKMCanvas(Canvas):
        def showPage(self):
            _stamp_footer(self, district, year)
            super().showPage()

        def save(self):
            _stamp_footer(self, district, year)
            super().save()

    return KKMCanvas


def _stamp_footer(canvas, district: str, year: int):
    text = FOOTER_TEMPLATE.format(district=district, year=year)
    canvas.saveState()
    canvas.setFont("Helvetica", 7)
    canvas.setFillColor(rl_colors.HexColor(KKM_MID_GRAY))
    canvas.drawCentredString(A4[0] / 2.0, 0.7 * cm, text)
    canvas.restoreState()


# --- section builders --------------------------------------------------------

def _pdf_section_cover(story, eda: dict, district: str, date_range: str,
                       cover_h, small):
    today  = date.today().strftime("%d %B %Y")
    # run_eda emits source_type/total_rows at the top level; the legacy
    # nested "summary" shape is kept only as a fallback.
    source = (eda.get("source_type")
              or eda.get("summary", {}).get("source_type", "N/A")).upper()
    rows   = eda.get("total_rows",
                     eda.get("summary", {}).get("total_rows", "N/A"))
    meta   = (
        f"Daerah / District: {district}  |  Sumber / Source: {source}"
        f"  |  Rekod / Records: {rows}"
        + (f"  |  Tempoh / Period: {date_range}" if date_range else "")
    )

    cover_tbl = Table(
        [
            [Paragraph("LAPORAN PEMAKANAN SmartDQC", cover_h)],
            [Paragraph("SMARTDQC NUTRITION REPORT",  cover_h)],
            [Paragraph(
                f'<font color="{KKM_TEAL_LIGHT}">{meta}<br/>{today}</font>',
                small,
            )],
        ],
        colWidths=[17 * cm],
    )
    cover_tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (-1, -1), rl_colors.HexColor(KKM_TEAL_DARK)),
        ("TOPPADDING",    (0, 0), (-1, -1), 14),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 14),
        ("LEFTPADDING",   (0, 0), (-1, -1), 16),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.append(cover_tbl)
    story.append(Spacer(1, 0.5 * cm))


def _pdf_section_exec_summary(story, narrative: dict, sec_label, h2, body):
    story.append(_section_bar_pdf(_sec("executive_summary", "en"),
                                  _sec("executive_summary", "bm"), sec_label))
    story.append(Spacer(1, 0.2 * cm))
    exec_sum = narrative.get("executive_summary", {})
    story.append(Paragraph("<b>Bahasa Malaysia</b>", h2))
    story.append(Paragraph(exec_sum.get("bm", "-"), body))
    story.append(Spacer(1, 0.25 * cm))
    story.append(Paragraph("<b>English</b>", h2))
    story.append(Paragraph(exec_sum.get("en", "-"), body))
    story.append(Spacer(1, 0.5 * cm))


def _pdf_section_quality(story, eda: dict, sec_label, h2, body, charts: set[str] | None = None):
    story.append(_section_bar_pdf(_sec("quality_overview", "en"),
                                  _sec("quality_overview", "bm"), sec_label))
    story.append(Spacer(1, 0.2 * cm))

    data = [["Metric / Metrik", "Value / Nilai"]]
    data.extend(_quality_overview_rows(eda))

    tbl = Table(data, colWidths=[11 * cm, 6 * cm])
    tbl.setStyle(TableStyle(_base_table_style()))
    story.append(tbl)

    if charts is None or "quality_bar" in charts:
        chart_png = chart_quality_bar(eda)
        if chart_png:
            img = Image(BytesIO(chart_png), width=14 * cm, height=7 * cm)
            story.append(Spacer(1, 0.3 * cm))
            story.append(img)

    story.append(Spacer(1, 0.5 * cm))


def _pdf_section_recommendations(story, narrative: dict, sec_label, h2, body, lang: str = "en"):
    recs = narrative.get("recommendations", [])
    if not recs:
        return
    story.append(_section_bar_pdf(_sec("recommendations", "en"),
                                  _sec("recommendations", "bm"), sec_label))
    story.append(Spacer(1, 0.2 * cm))
    for rec in recs[:5]:
        priority = rec.get("priority", "").upper()
        action = _action_for(rec, lang)
        story.append(Paragraph(
            f"<b>[{priority}] {action}</b>", h2))
        story.append(Paragraph(
            f"<b>BM:</b> {rec.get('bm', '-')}", body))
        story.append(Paragraph(
            f"<b>EN:</b> {rec.get('en', '-')}", body))
        story.append(Spacer(1, 0.25 * cm))
    story.append(Spacer(1, 0.3 * cm))


def _pdf_section_indicator_table(story, kpi_result: dict, sec_label, lang: str = "en",
                                 charts: set[str] | None = None):
    breakdown = kpi_result.get("by_daerah") or kpi_result.get("by_state") or []
    story.append(_section_bar_pdf(_sec("indicator_table", "en"),
                                  _sec("indicator_table", "bm"), sec_label))
    story.append(Spacer(1, 0.2 * cm))

    headers = NUTRITIONAL_TABLE_HEADERS[lang]
    data = [headers]
    for row in breakdown:
        rates = row.get("rates") or {}
        data.append([
            str(row.get("district") or row.get("state") or ""),
            str(row.get("n", "")),
            _fmt(rates.get("stunting")),
            _fmt(rates.get("wasting")),
            _fmt(rates.get("underweight")),
            _fmt(rates.get("overweight")),
        ])

    if len(data) > 1:
        tbl = Table(data, colWidths=[4.5*cm, 1.5*cm, 2.5*cm, 2.5*cm, 2.7*cm, 2.8*cm])
        tbl.setStyle(TableStyle(_base_table_style()))
        story.append(tbl)
    else:
        from reportlab.lib.styles import getSampleStyleSheet
        story.append(Paragraph("No district breakdown available.",
                               getSampleStyleSheet()["Normal"]))

    # Honour the per-chart filter; default (charts is None) embeds both.
    want_rates  = charts is None or "nutritional_rates" in charts
    want_target = charts is None or "kpi_vs_target" in charts
    chart2 = chart_nutritional_rates(kpi_result) if want_rates else None
    chart3 = chart_kpi_vs_target(kpi_result) if want_target else None
    if chart2 or chart3:
        story.append(Spacer(1, 0.3 * cm))
        chart_row = []
        if chart2:
            chart_row.append(Image(BytesIO(chart2), width=8.5 * cm, height=6.5 * cm))
        if chart3:
            chart_row.append(Image(BytesIO(chart3), width=8.5 * cm, height=6.5 * cm))
        if len(chart_row) == 2:
            chart_tbl = Table([chart_row], colWidths=[8.5 * cm, 8.5 * cm])
            chart_tbl.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP")]))
            story.append(chart_tbl)
        elif chart_row:
            story.append(chart_row[0])

    story.append(Spacer(1, 0.5 * cm))


def _pdf_section_methodology(story, sec_label, small):
    story.append(_section_bar_pdf(_sec("methodology", "en"),
                                  _sec("methodology", "bm"), sec_label))
    story.append(Spacer(1, 0.2 * cm))
    for line in METHODOLOGY_LINES:
        story.append(Paragraph(f"- {line}", small))
    story.append(Spacer(1, 0.3 * cm))


# --- public entry point ------------------------------------------------------

def build_pdf_bytes(
    eda_result: dict,
    narrative: dict,
    kpi_result: dict | None = None,
    district: str = "Malaysia",
    date_range: str = "",
    lang: str = "en",
    charts: set[str] | None = None,
) -> bytes:
    """Render the PDF report.

    `charts` controls which chart images are embedded:
      None         → default (every recommended chart)
      set of keys  → only those keys; known: quality_bar, nutritional_rates,
                     kpi_vs_target.
    """
    buf = BytesIO()
    footer_canvas = _make_footer_canvas(district, date.today().year)
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=2 * cm, bottomMargin=2 * cm,
    )

    cover_h, sec_label, h2, body, small = _pdf_styles()
    story: list = []

    _pdf_section_cover(story, eda_result, district, date_range, cover_h, small)
    _pdf_section_exec_summary(story, narrative, sec_label, h2, body)
    _pdf_section_quality(story, eda_result, sec_label, h2, body, charts=charts)
    _pdf_section_recommendations(story, narrative, sec_label, h2, body, lang)

    if kpi_result:
        _pdf_section_indicator_table(story, kpi_result, sec_label, lang, charts=charts)

    _pdf_section_methodology(story, sec_label, small)

    doc.build(story, canvasmaker=footer_canvas)
    return buf.getvalue()
